#!/usr/bin/env python3
"""
Enhanced HTML to PowerPoint Converter for Upflyover Pitch Deck
Converts the HTML pitch deck to a professional PowerPoint presentation with visual designs
"""

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

def create_upflyover_presentation():
    """Convert HTML pitch deck to PowerPoint presentation with designs"""
    
    # Read HTML file
    with open('upflyover-investor-pitch-deck.html', 'r', encoding='utf-8') as file:
        html_content = file.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Create presentation
    prs = Presentation()
    
    # Define colors (matching HTML theme)
    primary_color = RGBColor(30, 86, 86)  # #1e5656
    accent_color = RGBColor(255, 107, 53)  # #ff6b35
    light_gray = RGBColor(248, 249, 250)  # #f8f9fa
    white = RGBColor(255, 255, 255)
    green_color = RGBColor(40, 167, 69)  # #28a745
    
    # Extract slides
    slides = soup.find_all('div', class_='slide')
    
    for slide_div in slides:
        slide_content = slide_div.find('div', class_='slide-content')
        if not slide_content:
            continue
            
        # Add slide with background
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add gradient background
        add_gradient_background(slide, primary_color)
        
        # Extract slide number
        slide_number = slide_div.find('div', class_='slide-number')
        slide_num = slide_number.text if slide_number else ""
        
        # Add main content background shape
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.8), Inches(9), Inches(6.2)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = white
        content_bg.shadow.inherit = False
        
        # Extract title
        title = slide_content.find(['h1', 'h2'])
        if title:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1.2))
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = clean_text(title.get_text())
            title_p.font.size = Pt(36)
            title_p.font.color.rgb = primary_color
            title_p.font.bold = True
            title_p.alignment = PP_ALIGN.CENTER
        
        # Extract subtitle
        subtitle = slide_content.find('p', class_='subtitle')
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = clean_text(subtitle.get_text())
            subtitle_p.font.size = Pt(18)
            subtitle_p.font.color.rgb = RGBColor(102, 102, 102)
            subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Extract metrics
        metrics = slide_content.find_all('div', class_='metric-card')
        if metrics:
            y_pos = 3.8
            x_start = 1.2
            width = 2.2
            
            for i, metric in enumerate(metrics[:3]):  # Limit to 3 metrics per slide
                value = metric.find('div', class_='metric-value')
                label = metric.find('div', class_='metric-label')
                
                if value and label:
                    x_pos = x_start + (i * 2.6)
                    
                    # Add metric card background
                    metric_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos-0.1), Inches(y_pos-0.1), Inches(width), Inches(1.3)
                    )
                    metric_bg.fill.solid()
                    metric_bg.fill.fore_color.rgb = light_gray
                    metric_bg.line.color.rgb = primary_color
                    metric_bg.line.width = Pt(3)
                    
                    # Add metric text
                    metric_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(width-0.2), Inches(1.1))
                    metric_frame = metric_box.text_frame
                    
                    # Value
                    value_p = metric_frame.paragraphs[0]
                    value_p.text = clean_text(value.get_text())
                    value_p.font.size = Pt(28)
                    value_p.font.color.rgb = primary_color
                    value_p.font.bold = True
                    value_p.alignment = PP_ALIGN.CENTER
                    
                    # Label
                    label_p = metric_frame.add_paragraph()
                    label_p.text = clean_text(label.get_text())
                    label_p.font.size = Pt(12)
                    label_p.font.color.rgb = RGBColor(102, 102, 102)
                    label_p.alignment = PP_ALIGN.CENTER
        
        # Extract bullet points
        lists = slide_content.find_all('ul')
        if lists:
            y_pos = 4.8 if not metrics else 5.8
            
            for ul in lists:
                list_items = ul.find_all('li')
                if list_items:
                    bullet_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(7.6), Inches(2.5))
                    bullet_frame = bullet_box.text_frame
                    
                    for i, li in enumerate(list_items[:4]):  # Limit to 4 items
                        if i == 0:
                            p = bullet_frame.paragraphs[0]
                        else:
                            p = bullet_frame.add_paragraph()
                        
                        p.text = f"• {clean_text(li.get_text())}"
                        p.font.size = Pt(14)
                        p.font.color.rgb = RGBColor(51, 51, 51)
                        p.space_after = Pt(8)
        
        # Extract highlight boxes
        highlights = slide_content.find_all('div', class_='highlight')
        for highlight in highlights:
            amount = highlight.find('div', class_='amount')
            if amount:
                # Add highlight background shape
                highlight_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3.2), Inches(6), Inches(2.2)
                )
                highlight_bg.fill.solid()
                highlight_bg.fill.fore_color.rgb = accent_color
                highlight_bg.shadow.inherit = False
                
                # Add highlight text
                highlight_box = slide.shapes.add_textbox(Inches(2.2), Inches(3.4), Inches(5.6), Inches(1.8))
                highlight_frame = highlight_box.text_frame
                
                # Title
                h3 = highlight.find('h3')
                if h3:
                    title_p = highlight_frame.paragraphs[0]
                    title_p.text = clean_text(h3.get_text())
                    title_p.font.size = Pt(20)
                    title_p.font.bold = True
                    title_p.font.color.rgb = white
                    title_p.alignment = PP_ALIGN.CENTER
                
                # Amount
                amount_p = highlight_frame.add_paragraph()
                amount_p.text = clean_text(amount.get_text())
                amount_p.font.size = Pt(36)
                amount_p.font.bold = True
                amount_p.font.color.rgb = white
                amount_p.alignment = PP_ALIGN.CENTER
                
                # Description
                desc_p = highlight_frame.add_paragraph()
                desc_text = highlight.get_text().replace(h3.get_text(), '').replace(amount.get_text(), '')
                desc_p.text = clean_text(desc_text)[:80] + '...' if len(desc_text) > 80 else clean_text(desc_text)
                desc_p.font.size = Pt(12)
                desc_p.font.color.rgb = white
                desc_p.alignment = PP_ALIGN.CENTER
        
        # Extract competitive advantage boxes
        comp_adv = slide_content.find_all('div', class_='competitive-advantage')
        for comp in comp_adv:
            # Add green background shape
            comp_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.5), Inches(7), Inches(3)
            )
            comp_bg.fill.solid()
            comp_bg.fill.fore_color.rgb = green_color
            
            # Add text
            comp_box = slide.shapes.add_textbox(Inches(1.7), Inches(3.7), Inches(6.6), Inches(2.6))
            comp_frame = comp_box.text_frame
            
            comp_title = comp.find('h3')
            if comp_title:
                title_p = comp_frame.paragraphs[0]
                title_p.text = clean_text(comp_title.get_text())
                title_p.font.size = Pt(22)
                title_p.font.bold = True
                title_p.font.color.rgb = white
                title_p.alignment = PP_ALIGN.CENTER
        
        # Add slide number with background
        if slide_num:
            # Number background
            num_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(0.3), Inches(1.2), Inches(0.4)
            )
            num_bg.fill.solid()
            num_bg.fill.fore_color.rgb = primary_color
            
            # Number text
            num_box = slide.shapes.add_textbox(Inches(8.3), Inches(0.35), Inches(1), Inches(0.3))
            num_frame = num_box.text_frame
            num_p = num_frame.paragraphs[0]
            num_p.text = slide_num
            num_p.font.size = Pt(10)
            num_p.font.color.rgb = white
            num_p.font.bold = True
            num_p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Upflyover_Investor_Pitch_Deck_Designed.pptx')
    print("✅ Designed PowerPoint presentation created: Upflyover_Investor_Pitch_Deck_Designed.pptx")

def add_gradient_background(slide, primary_color):
    """Add gradient background to slide"""
    # Add background rectangle
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(45, 125, 125)  # Lighter version of primary
    bg_shape.line.fill.background()
    
    # Move to back
    slide.shapes._spTree.remove(bg_shape._element)
    slide.shapes._spTree.insert(2, bg_shape._element)

def clean_text(text):
    """Clean text by removing extra whitespace and emojis for PowerPoint compatibility"""
    # Remove emojis and special characters that might cause issues
    text = re.sub(r'[^\w\s\-\.\,\:\;\!\?\$\%\(\)\/\+\=]', '', text)
    # Clean whitespace
    return ' '.join(text.split())

def install_requirements():
    """Install required packages"""
    import subprocess
    import sys
    
    packages = ['python-pptx', 'beautifulsoup4']
    
    for package in packages:
        try:
            __import__(package.replace('-', '_'))
        except ImportError:
            print(f"Installing {package}...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", package])

if __name__ == "__main__":
    try:
        install_requirements()
        create_upflyover_presentation()
    except Exception as e:
        print(f"❌ Error: {e}")
        print("Make sure you have the HTML file in the same directory and required packages installed.")