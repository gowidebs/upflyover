#!/usr/bin/env python3
"""
HTML to PowerPoint Converter for Upflyover Pitch Deck
Converts the HTML pitch deck to a professional PowerPoint presentation
"""

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

def create_upflyover_presentation():
    """Convert HTML pitch deck to PowerPoint presentation"""
    
    # Read HTML file
    with open('upflyover-investor-pitch-deck.html', 'r', encoding='utf-8') as file:
        html_content = file.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Create presentation
    prs = Presentation()
    
    # Define colors (matching HTML theme)
    primary_color = RGBColor(30, 86, 86)  # #1e5656
    accent_color = RGBColor(255, 107, 53)  # #ff6b35
    
    # Extract slides
    slides = soup.find_all('div', class_='slide')
    
    for slide_div in slides:
        slide_content = slide_div.find('div', class_='slide-content')
        if not slide_content:
            continue
            
        # Add slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Extract slide number
        slide_number = slide_div.find('div', class_='slide-number')
        slide_num = slide_number.text if slide_number else ""
        
        # Extract title
        title = slide_content.find(['h1', 'h2'])
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = clean_text(title.get_text())
            title_p.font.size = Pt(36)
            title_p.font.color.rgb = primary_color
            title_p.font.bold = True
            title_p.alignment = PP_ALIGN.CENTER
        
        # Extract subtitle
        subtitle = slide_content.find('p', class_='subtitle')
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = clean_text(subtitle.get_text())
            subtitle_p.font.size = Pt(18)
            subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Extract metrics
        metrics = slide_content.find_all('div', class_='metric-card')
        if metrics:
            y_pos = 3.5
            x_start = 1
            width = 2.5
            
            for i, metric in enumerate(metrics[:3]):  # Limit to 3 metrics per slide
                value = metric.find('div', class_='metric-value')
                label = metric.find('div', class_='metric-label')
                
                if value and label:
                    x_pos = x_start + (i * 2.8)
                    
                    # Add metric box
                    metric_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(width), Inches(1.5))
                    metric_frame = metric_box.text_frame
                    
                    # Value
                    value_p = metric_frame.paragraphs[0]
                    value_p.text = clean_text(value.get_text())
                    value_p.font.size = Pt(32)
                    value_p.font.color.rgb = primary_color
                    value_p.font.bold = True
                    value_p.alignment = PP_ALIGN.CENTER
                    
                    # Label
                    label_p = metric_frame.add_paragraph()
                    label_p.text = clean_text(label.get_text())
                    label_p.font.size = Pt(14)
                    label_p.alignment = PP_ALIGN.CENTER
        
        # Extract bullet points
        lists = slide_content.find_all('ul')
        if lists:
            y_pos = 4.5 if not metrics else 5.5
            
            for ul in lists:
                list_items = ul.find_all('li')
                if list_items:
                    bullet_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(3))
                    bullet_frame = bullet_box.text_frame
                    
                    for i, li in enumerate(list_items[:5]):  # Limit to 5 items
                        if i == 0:
                            p = bullet_frame.paragraphs[0]
                        else:
                            p = bullet_frame.add_paragraph()
                        
                        p.text = f"• {clean_text(li.get_text())}"
                        p.font.size = Pt(16)
                        p.space_after = Pt(6)
        
        # Extract highlight boxes
        highlights = slide_content.find_all('div', class_='highlight')
        for highlight in highlights:
            amount = highlight.find('div', class_='amount')
            if amount:
                highlight_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
                highlight_frame = highlight_box.text_frame
                
                # Title
                h3 = highlight.find('h3')
                if h3:
                    title_p = highlight_frame.paragraphs[0]
                    title_p.text = clean_text(h3.get_text())
                    title_p.font.size = Pt(24)
                    title_p.font.bold = True
                    title_p.alignment = PP_ALIGN.CENTER
                
                # Amount
                amount_p = highlight_frame.add_paragraph()
                amount_p.text = clean_text(amount.get_text())
                amount_p.font.size = Pt(48)
                amount_p.font.bold = True
                amount_p.font.color.rgb = accent_color
                amount_p.alignment = PP_ALIGN.CENTER
        
        # Add slide number
        if slide_num:
            num_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.2), Inches(1), Inches(0.5))
            num_frame = num_box.text_frame
            num_p = num_frame.paragraphs[0]
            num_p.text = slide_num
            num_p.font.size = Pt(12)
            num_p.alignment = PP_ALIGN.RIGHT
    
    # Save presentation
    prs.save('Upflyover_Investor_Pitch_Deck.pptx')
    print("✅ PowerPoint presentation created: Upflyover_Investor_Pitch_Deck.pptx")

def clean_text(text):
    """Clean text by removing extra whitespace and emojis for PowerPoint compatibility"""
    # Remove emojis and special characters that might cause issues
    text = re.sub(r'[^\w\s\-\.\,\:\;\!\?\$\%\(\)\/\+\=]', '', text)
    # Clean whitespace
    return ' '.join(text.split())

def install_requirements():
    """Install required packages"""
    import subprocess
    import sys
    
    packages = ['python-pptx', 'beautifulsoup4']
    
    for package in packages:
        try:
            __import__(package.replace('-', '_'))
        except ImportError:
            print(f"Installing {package}...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", package])

if __name__ == "__main__":
    try:
        install_requirements()
        create_upflyover_presentation()
    except Exception as e:
        print(f"❌ Error: {e}")
        print("Make sure you have the HTML file in the same directory and required packages installed.")